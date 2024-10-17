# Módulo de RAG


import sys  # Importa o módulo sys para acessar os argumentos da linha de comando
import docx # Importa o módulo docx para manipulação de arquivos Word
import PyPDF2   # Importa o módulo PyPDF2 para manipulação de arquivos PDF
from pptx import Presentation   # Importa o módulo Presentation do pacote pptx para manipulação de arquivos PowerPoint
from os import listdir  # Importa as funções listdir, isfile, join e isdir dos módulos os e os.path para manipulação de diretórios e arquivos
from os.path import isfile, join, isdir
from langchain_text_splitters import TokenTextSplitter # Importa a classe TokenTextSplitter do pacote langchain_text_splitters para a divisão do texto em tokens
from langchain_huggingface import HuggingFaceEmbeddings # Importa a classe HuggingFaceEmbeddings do pacote langchain_huggingface para criar as embeddings
from qdrant_client import QdrantClient #QdrantVectorStore  # Importa as classes QdrantClient, Distance e VectorParams do pacote qdrant_client
from qdrant_client.models import Distance, VectorParams
from langchain_qdrant import Qdrant # Importa a classe Qdrant do pacote langchain_qdrant para criar uma instância do Qdrant e enviar os dados para o banco vetorial


#----------------------------------------------------------------------------------------------------------------------------------
# Define a função que lista todos os arquivos em um diretório, incluindo os de subdiretórios
def lista_arquivos(dir):
        
    arquivo_list = []   # Inicializa uma lista vazia para armazenar os caminhos dos arquivos
        
    for f in listdir(dir):  # Itera sobre todos os arquivos e diretórios no diretório especificado
        if isfile(join(dir, f)):    # Se for um arquivo, adiciona à lista
            arquivo_list.append(join(dir, f))
                
        elif isdir(join(dir, f)):   # Se for um diretório, chama a função recursivamente e adiciona os resultados à lista
            arquivo_list += lista_arquivos(join(dir, f))
    
    return arquivo_list # Retorna a lista de arquivos

#----------------------------------------------------------------------------------------------------------------------------------
# Define a função que carrega o texto de um arquivo Word
def carrega_texto_word(arquivoname):
        
    doc = docx.Document(arquivoname)    # Abre o arquivo Word
    fullText = [para.text for para in doc.paragraphs]   # Extrai o texto de cada parágrafo e adiciona à lista
    return '\n'.join(fullText)  # Junta todos os textos em uma única string separada por quebras de linha

#----------------------------------------------------------------------------------------------------------------------------------
# Define a função que carrega o texto de um arquivo PowerPoint
def carrega_texto_pptx(arquivoname):
    
    prs = Presentation(arquivoname) # Abre o arquivo PowerPoint
    fullText = []   # Inicializa uma lista vazia para armazenar os textos
    for slide in prs.slides:    # Itera sobre todos os slides  
        for shape in slide.shapes:  # Itera sobre todas as formas no slide
            if hasattr(shape, "text"):  # Se a forma tiver o atributo "text", adiciona o texto à lista
                fullText.append(shape.text)
    
    return '\n'.join(fullText)  # Junta todos os textos em uma única string separada por quebras de linha

#----------------------------------------------------------------------------------------------------------------------------------
# Define a função principal para indexação dos documentos
def main_indexing(mypath):

    model_name = "sentence-transformers/msmarco-bert-base-dot-v5"   # Define o nome do modelo a ser usado para criar as embeddings
    model_kwargs = {'device': 'cpu'}    # Define as configurações do modelo
    encode_kwargs = {'normalize_embeddings': True}  # Define as configurações de codificação
    hf = HuggingFaceEmbeddings(model_name = model_name,         # Inicializa a classe de embeddings do HuggingFace
                               model_kwargs = model_kwargs,
                               encode_kwargs = encode_kwargs)
    client = QdrantClient("http://localhost:6333")  # Inicializa o cliente Qdrant
    collection_name = "vectordb"     # Define o nome da coleção de embeddings
    
    if client.collection_exists(collection_name):   # Se a coleção já existir, exclui
        client.delete_collection(collection_name)
    client.create_collection(collection_name,       # Cria uma nova coleção com parâmetros especificados
                             vectors_config = VectorParams(size = 768, distance = Distance.DOT))
    
    qdrant = Qdrant(client, collection_name, hf)        # Inicializa a instância Qdrant
    # Use a nova classe QdrantVectorStore em vez da deprecada Qdrant
    #qdrant = QdrantVectorStore(client=client, collection_name=collection_name, embedding_function=hf)
    print("\nIndexando os documentos...\n")     # Imprime mensagem informando que a indexação dos documentos está iniciando

    arquivos = lista_arquivos(mypath)     # Obtém a lista de todos os arquivos no diretório especificado
    print(arquivos)
    # Itera sobre cada arquivo na lista
    for arquivo in arquivos:
        
        try:
            
            arquivo_content = ""        # Inicializa uma string vazia para armazenar o conteúdo do arquivo

            if arquivo.endswith(".pdf"):    # Verifica se o arquivo é um PDF
                print("Indexando: " + arquivo)
                reader = PyPDF2.PdfReader(arquivo)
                
                for page in reader.pages:
                    arquivo_content += " " + page.extract_text()

            elif arquivo.endswith(".txt"):  # Verifica se o arquivo é um texto simples
                print("Indexando: " + arquivo)
                with open(arquivo, 'r') as f:
                    arquivo_content = f.read()
            
            elif arquivo.endswith(".docx"):     # Verifica se o arquivo é um Word
                print("Indexando: " + arquivo)
                arquivo_content = carrega_texto_word(arquivo)
            
            
            elif arquivo.endswith(".pptx"):     # Verifica se o arquivo é um PowerPoint
                print("Indexando: " + arquivo)
                arquivo_content = carrega_texto_pptx(arquivo)

            else:
                continue    # Se o arquivo não for de um formato suportado, continua para o próximo arquivo

            
            text_splitter = TokenTextSplitter(chunk_size = 500, chunk_overlap = 50) # Inicializa o divisor de texto com tamanho de chunk e sobreposição especificados
            textos = text_splitter.split_text(arquivo_content)      # Divide o conteúdo do arquivo em chunks de texto
            metadata = [{"path": arquivo} for _ in textos]      # Cria metadados para cada chunk de texto - isso permite que o LLM cite a referência
            qdrant.add_texts(textos, metadatas = metadata)      # Adiciona os textos e seus metadatas ao Qdrant

        except Exception as e:
            print(f"O processo falhou para o arquivo {arquivo}: {e}")   # Se ocorrer um erro, imprime uma mensagem de erro
 
    print("\nIndexação Concluída!\n")   # Imprime mensagem informando que a indexação foi concluída

#----------------------------------------------------------------------------------------------------------------------------------
if __name__ == "__main__":      # Verifica se o script está sendo executado diretamente

    arguments = sys.argv        # Obtém os argumentos da linha de comando
    
    if len(arguments) > 1:      # Verifica se foi fornecido um caminho de diretório
        main_indexing(arguments[1])
    else:
        print("Fornecer um caminho para a pasta com documentos para indexar.") # Se não, imprime uma mensagem de erro